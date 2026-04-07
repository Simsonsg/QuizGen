"""
Parse PDF, PPTX, or plain text files into raw text.
"""

import pdfplumber
from pptx import Presentation
from pathlib import Path


def parse_file(path: str) -> str:
    """
    Parse a document file into a single raw text string.
    Supports .pdf, .pptx, and .txt files.
    """
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {path}")

    suffix = p.suffix.lower()

    if suffix == ".pdf":
        return _parse_pdf(path)
    elif suffix == ".pptx":
        return _parse_pptx(path)
    elif suffix == ".txt":
        return p.read_text(encoding="utf-8")
    else:
        raise ValueError(f"Unsupported file type: {suffix}. Use .pdf, .pptx, or .txt")


def _parse_pdf(path: str) -> str:
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    return "\n\n".join(pages)


def _parse_pptx(path: str) -> str:
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_text.append(line)
        if slide_text:
            slides.append("\n".join(slide_text))
    return "\n\n".join(slides)
